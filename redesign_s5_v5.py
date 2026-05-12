"""
Redesign S5 only — replace the v4 "three contradictions" content with a clean
3-column "what is phase segregation" explainer using slide-4 layout style.

Approach: clear S5 → clone slide-4's shape XML into S5 → replace text by
position-based slot matching. Same layout as S4 (3 columns with FIG/PROBLEM/
DIRECTIONS/METRIC), but content is definitional, not analytical.

Formulas are intentionally NOT on the slide — they live in formulas.md and
the user screenshots them in. The 3 FIG slots are placeholders for canonical
figures (Hoke 2015, Yu 2026, Chen 2020) the user copies in manually.
"""
from __future__ import annotations

import copy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

SRC = Path("oral exam_260423_v4_SLPCI.pptx")
OUT = Path("oral exam_260423_v5_SLPCI.pptx")

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


# ============================================================
# NEW S5 CONTENT — definitional, with 3 figure placeholders
# ============================================================
TITLE = "4. Photoinduced halide phase segregation — what it is, why it matters"

SUBTITLE = (
    "Under continuous illumination, mixed-halide wide-bandgap perovskites spontaneously demix "
    "into iodide-rich (low-Eg) and bromide-rich (high-Eg) domains within minutes — capping the "
    "V_OC of wide-bandgap photovoltaics and limiting all-perovskite tandem efficiency."
)

PAGE = "0 5"

COLS = [
    # ---------- COLUMN 1 ----------
    dict(
        num="I.",
        sub="THE PHENOMENON",
        fig_label="FIG. 1",
        fig_caption="Hoke 2015, Fig. 2 (a-c).  MAPb(Br0.4I0.6)3 PL evolution under 457 nm @ 15 mW/cm², 0-45 s in 5 s steps; (b) intensity dependence; (c) dark recovery cycles.",
        problem=(
            "Under continuous 50-100 mW/cm² illumination the initial alloy-bandgap "
            "emission (e.g., 1.85 eV in MAPb(Br0.4I0.6)3) decays within 60-100 s. "
            "A new red-shifted PL peak at ~1.65 eV grows in."
        ),
        directions=(
            "The same film recovers fully in dark within 5-10 min and "
            "re-segregates on the next light cycle.\n"
            "Reversibility distinguishes ionic redistribution from chemical "
            "degradation — the canonical 'Hoke effect', first reported in 2015."
        ),
        metric="ΔE ≈ 0.20 eV",
        metric_caption="REVERSIBLE   ·   t_onset ≈ 60 s",
    ),

    # ---------- COLUMN 2 ----------
    dict(
        num="II.",
        sub="THE MECHANISM",
        fig_label="FIG. 2",
        fig_caption="Yu 2026, Fig. 3b (originally Draguta 2017, ref 89).  Filled dark circles = photogenerated e-h pairs; empty circles = pairs that induced segregation; red = I-rich domains; white = Br-rich matrix.",
        problem=(
            "Photogenerated electron-hole pairs polarize the soft halide lattice "
            "via polaron formation. The local free-energy gain drives I⁻ ions "
            "into ~100 nm I-rich domains; Br⁻ accumulates in the surrounding "
            "matrix."
        ),
        directions=(
            "Carriers then funnel into the low-Eg I-rich domains — emission "
            "and quasi-Fermi splitting are dominated by these inclusions "
            "even when their volume fraction is < 0.1 %.\n"
            "Halide vacancy hopping (E_a ≈ 0.3 eV) is the rate-limiting step."
        ),
        metric="E_a ≈ 0.3 eV",
        metric_caption="I-RICH DOMAIN ~ 100 nm",
    ),

    # ---------- COLUMN 3 ----------
    dict(
        num="III.",
        sub="WHY IT MATTERS",
        fig_label="FIG. 3",
        fig_caption="Fang 2024, Fig. 1e (originally Wiley-VCH 2023, ref 26).  V_OC deficit vs bandgap for WBG PSCs: > 1.72 eV cells systematically lose > 0.5 V.",
        problem=(
            "Once I-rich domains form, carriers thermalize into them and emit "
            "at 1.65 eV — pinning quasi-Fermi splitting far below the alloy "
            "bandgap. WBG ~1.8 eV cells, needed as tandem top cells, lose "
            "> 0.5 V of V_OC."
        ),
        directions=(
            "All-perovskite tandems are direct victims — top-cell V_OC loss "
            "propagates one-to-one into tandem PCE.\n"
            "Suppressing or quantitatively understanding phase segregation "
            "is the bottleneck to > 30 % certified tandem efficiency."
        ),
        metric="V_OC loss > 0.5 V",
        metric_caption="WBG > 1.72 eV  ·  TANDEM PCE CAP",
    ),
]


# ============================================================
# Helpers
# ============================================================
def set_tf_text(tf, new_text: str) -> None:
    """Replace tf contents (multi-line supported) preserving first run's font."""
    template_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            template_run = r
            break
        if template_run is not None:
            break
    txBody = tf._txBody
    for p in list(txBody.findall("a:p", NSMAP)):
        txBody.remove(p)
    for line in new_text.split("\n"):
        p_el = etree.SubElement(txBody, f"{A_NS}p")
        if template_run is not None:
            new_run = copy.deepcopy(template_run._r)
            for t in new_run.findall(f"{A_NS}t"):
                new_run.remove(t)
            t_el = etree.SubElement(new_run, f"{A_NS}t")
            t_el.text = line
            p_el.append(new_run)
        else:
            r_el = etree.SubElement(p_el, f"{A_NS}r")
            t_el = etree.SubElement(r_el, f"{A_NS}t")
            t_el.text = line


def first_text(sh) -> str:
    if not sh.has_text_frame:
        return ""
    return " | ".join((p.text or "").strip()
                      for p in sh.text_frame.paragraphs
                      if (p.text or "").strip())


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)


def clone_shapes(src_slide, dst_slide) -> None:
    """Deep-copy all top-level shape XML from src into dst."""
    src_spTree = src_slide.shapes._spTree
    dst_spTree = dst_slide.shapes._spTree
    for child in src_spTree:
        tag = etree.QName(child.tag).localname
        if tag in ("nvGrpSpPr", "grpSpPr"):
            continue
        dst_spTree.append(copy.deepcopy(child))


# ============================================================
# Slot mapping (position-based) — adapted from v2 fill_slide()
# ============================================================
def fill_s5(slide, content) -> None:
    title_text = content["title"]
    sub_text = content["subtitle"]
    page_text = content["page"]
    cols = content["cols"]

    col_xranges = [(0.0, 4.5), (4.5, 8.7), (8.7, 13.4)]

    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        try:
            lx, ly = Emu(sh.left).inches, Emu(sh.top).inches
            w = Emu(sh.width).inches
        except Exception:
            continue
        txt = first_text(sh)
        if not txt:
            continue

        # Title (top-left, narrow row)
        if ly < 0.3 and lx < 0.5 and w > 5:
            set_tf_text(sh.text_frame, title_text)
            continue
        # Subtitle band
        if 0.55 < ly < 0.75 and w > 10:
            set_tf_text(sh.text_frame, sub_text)
            continue
        # Page number
        if ly > 7.2 and lx > 12.5:
            set_tf_text(sh.text_frame, page_text)
            continue

        # Identify column
        col_idx = None
        for ci, (xmin, xmax) in enumerate(col_xranges):
            if xmin <= lx < xmax:
                col_idx = ci
                break
        if col_idx is None:
            continue
        col = cols[col_idx]

        # I./II./III. roman numeral (narrow, T~1.28)
        if 1.20 < ly < 1.32 and w < 0.6 and txt in ("I.", "II.", "III."):
            set_tf_text(sh.text_frame, col["num"])
            continue
        # Column sub-title (Stability/Upscaling/Efficiency)
        if 1.30 < ly < 1.40 and 0.6 < w < 1.5 and txt in ("Stability", "Upscaling", "Efficiency"):
            set_tf_text(sh.text_frame, col["sub"])
            continue
        # FIG. n label
        if 1.82 < ly < 1.92 and txt.startswith("FIG."):
            set_tf_text(sh.text_frame, col["fig_label"])
            continue
        # Suggested fig caption
        if 2.80 < ly < 2.95 and txt.startswith("Suggested:"):
            set_tf_text(sh.text_frame, col["fig_caption"])
            continue
        # PROBLEM body (T~4.44)
        if 4.35 < ly < 4.55 and w > 3.0 and not txt.startswith("PROBLEM"):
            set_tf_text(sh.text_frame, col["problem"])
            continue
        # DIRECTIONS body (T~5.32, multi-line)
        if 5.25 < ly < 5.45 and w > 3.0:
            set_tf_text(sh.text_frame, col["directions"])
            continue
        # Bottom metric caption FIRST (matched by keyword; value/caption Y ranges overlap)
        if 6.20 < ly < 6.50 and any(kw in txt for kw in (
            "OPERATIONAL TARGET", "MODULE-SCALE", "SINGLE-JUNCTION", "TARGET", "RECORD"
        )):
            set_tf_text(sh.text_frame, col["metric_caption"])
            continue
        # Bottom metric value (whatever narrow text remains at T~6.17-6.35)
        if 6.10 < ly < 6.35 and w < 1.2:
            set_tf_text(sh.text_frame, col["metric"])
            continue


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation(str(SRC))
    print(f"Opened: {SRC.name}  ({len(prs.slides)} slides)")

    s4 = prs.slides[3]
    s5 = prs.slides[4]

    print("Clearing current S5 ...")
    clear_slide(s5)

    print("Cloning slide-4 layout into S5 ...")
    clone_shapes(s4, s5)

    print("Filling new S5 with phase-segregation explainer content ...")
    fill_s5(s5, dict(title=TITLE, subtitle=SUBTITLE, page=PAGE, cols=COLS))

    prs.save(str(OUT))
    print(f"\nSaved: {OUT}")


if __name__ == "__main__":
    main()
