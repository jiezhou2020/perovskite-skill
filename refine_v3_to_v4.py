"""
Refine v3 → v4: keep all layouts unchanged, only tighten English text.

Targets:
  S2  — translate Chinese background bullets to English, polish work-content
  S3  — skip per user instruction
  S5–S10  — tighten subtitles and card bodies to PPT presentation style:
            1-sentence subtitles, 2-clause body sentences, concrete numbers,
            visible citations.

Strategy: open v3 pptx, locate target text frames by content prefix,
replace text while preserving the first run's font formatting.
"""
from __future__ import annotations

import copy
from pathlib import Path

from lxml import etree
from pptx import Presentation

SRC = Path("oral exam_260423_v3_SLPCI.pptx")
OUT = Path("oral exam_260423_v4_SLPCI.pptx")

A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def get_full_text(tf) -> str:
    return "\n".join((p.text or "") for p in tf.paragraphs)


def set_tf_text(tf, new_text: str) -> None:
    """Replace tf contents (multi-line supported) preserving first run's font."""
    # Capture first run as formatting template
    template_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            template_run = r
            break
        if template_run is not None:
            break

    txBody = tf._txBody
    # Remove all existing <a:p> children
    for p in list(txBody.findall("a:p", NSMAP)):
        txBody.remove(p)

    # Write new paragraphs
    for line in new_text.split("\n"):
        p_el = etree.SubElement(txBody, f"{A_NS}p")
        if template_run is not None:
            new_run = copy.deepcopy(template_run._r)
            for t in new_run.findall(f"{A_NS}t"):
                new_run.remove(t)
            t_el = etree.SubElement(new_run, f"{A_NS}t")
            t_el.text = line
            p_el.append(new_run)
        else:
            r_el = etree.SubElement(p_el, f"{A_NS}r")
            t_el = etree.SubElement(r_el, f"{A_NS}t")
            t_el.text = line


def replace_in_slide(slide, match_substr: str, new_text: str, label: str = "") -> bool:
    """Find the first text frame whose text contains `match_substr` and replace it."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        full = get_full_text(sh.text_frame)
        if match_substr in full:
            set_tf_text(sh.text_frame, new_text)
            if label:
                print(f"   [{label}] replaced  '{match_substr[:50]}...'  →  '{new_text[:50]}...'")
            return True
    print(f"   [{label}]  NO MATCH for prefix '{match_substr[:60]}'")
    return False


# ============================================================
# S2 — Background: Chinese → polished English
# ============================================================
S2_REPLACEMENTS = [
    # — Waseda block —
    ("研究方向：钙钛矿薄膜沉积技术",
     "Research focus:  large-area perovskite deposition  ·  EIP apparatus  ·  perovskite solar cells"),

    ("先行研究调研：①钙钛矿薄膜大面积沉积技术",
     "Literature review:  ① large-area perovskite deposition    ② electrostatic inkjet printing (EIP)"),

    ("设备改善方案设计",
     "EIP apparatus improvement proposals  ·  Taylor-cone control + Coulomb fission"),

    ("调研文献整理",
     "Reference compilation  +  process-flow documentation"),

    ("M2: 实验室实验",
     "M2:  in-lab experiments  +  thesis writing"),

    ("实验装置搭建",
     "Built EIP apparatus from scratch  (chamber + nozzle + HV supply)"),

    ("设备参数调优",
     "Optimized voltage  /  flow-rate  /  nozzle  /  substrate parameters"),

    ("综述论文 and 毕业论文",
     "Review paper  +  Master's thesis  (large-area EIP for perovskite films)"),

    ("论文1",
     "1 conference paper  (CSJ Tokyo 2023)"),

    ("学会发表1",
     "1 conference presentation"),

    ("2023.04-至今",
     "2023.04 – present"),
]


# ============================================================
# S5 — Coupled-state problem: tighten
# ============================================================
S5_REPLACEMENTS = [
    # subtitle (1 sentence)
    ("Three contradictions in the literature are not measurement noise.",
     "Three contradictions in the literature are not measurement noise — they are evidence that each technique reports a different projection of one shared physical state."),

    # card 1 body — concrete V_D numbers + tighter ending
    ("Minority I-rich domains capture carriers from a diffusion volume V_D and dominate emission",
     "Carriers funnel into low-E_g I-rich domains over V_D ≈ 100–1000 nm, so emission is dominated even at φ_I ≈ 10⁻³.  PL red-shift means carriers found a low-E_g site — not how much phase formed."),

    # card 2 body — keep equation, tighten ending
    ("V_CPD = W₀(χ_Br_surf) − V_bb(N_defect)",
     "V_CPD mixes  W₀ , V_bb , V_SPV , V_ion , V_trap.  A single static scan cannot isolate any one component — yet 'excess CPD ⇒ local electric field' remains the field's default reading (Qu 2026)."),

    # card 3 body — tighter
    ("X-ray dose (XPS) and sputter (TOF-SIMS) move halides.",
     "X-ray dose (XPS) and ion sputter (TOF-SIMS) themselves drive halide migration.  Reported 'Br-rich surfaces' may be probe artifacts — the Fang 2024 ↔ Navid 2026 directional disagreement is the canonical symptom."),
]


# ============================================================
# S6 — Central hypothesis: tighten subtitle, polish claims
# ============================================================
S6_REPLACEMENTS = [
    ("Phase segregation is governed by five coupled latent states. Br richness is neither necessary nor sufficient",
     "Five coupled latent states govern phase segregation.  Br richness alone is neither necessary nor sufficient — its phenotype depends on whether the other four co-move with it."),

    # Hα claim — already tight; just polish punctuation
    ("Br-rich surface stabilizes only when it simultaneously reduces N_defect, smooths φ_local, and releases local strain ε.",
     "Br-rich surface stabilizes only when it simultaneously reduces N_defect , smooths φ_local AND releases ε."),

    # Hβ claim — already tight
    ("Local CPD anomaly is a trigger only if it precedes the PL hotspot in time AND co-locates in the same ROI.",
     "Local CPD anomaly triggers LHS only if it precedes the PL hotspot in time AND co-locates in the same ROI."),

    # Hγ claim — already tight
    ("Reversible LHS and irreversible degradation share PL signatures — separable only by joint posterior across optical + chemical + gas channels.",
     "Reversible LHS and irreversible degradation share PL signatures — separable only by the joint posterior across optical , chemical and gas channels."),
]


# ============================================================
# S7 — SL-PCI method
# ============================================================
S7_REPLACEMENTS = [
    ("State-Locked Physics-Informed Correlative Inference — write a defined state",
     "State-Locked Physics-Informed Correlative Inference:  write a defined state  →  freeze it for endpoint probes  →  infer one shared latent state Z with explicit artifact channels."),

    # Step 4 body — tighten
    ("Bayesian inference returns p(Z, θ | y) with covariance. Jacobian SVD identifies ill-conditioned directions before any sample is made.",
     "Bayesian inference returns p(Z, θ | y) with covariance.  Jacobian SVD identifies which Z directions are identifiable before any sample is fabricated — output: mechanism posteriors + Bayes factors."),
]


# ============================================================
# S8 — Forward models
# ============================================================
S8_REPLACEMENTS = [
    ("Listing H_m(Z; θ_m) explicitly makes the inference reproducible. Each technique's primary observable, governing equation and resolved Z components are laid out before any sample is fabricated.",
     "Listing H_m(Z; θ_m) explicitly makes the inference reproducible.  Each technique's observable, governing equation and resolved Z components are fixed before any sample is fabricated."),

    ("Every probe is over-determined when combined — no single instrument resolves Z, but their joint likelihood does.",
     "No single instrument resolves Z — but the joint likelihood across all H_m does, and over-determination makes the inference robust to any one probe's failure."),
]


# ============================================================
# S9 — Identifiability budget
# ============================================================
S9_REPLACEMENTS = [
    ("Jacobian J_mj = ∂H_m / ∂Z_j is built from the forward models before any sample exists. Its SVD reveals which mechanism directions a probe set constrains. Aim 1→3 adds probes that close specific ill-conditioned directions; condition number κ tracks progress.",
     "Jacobian J_mj = ∂H_m / ∂Z_j is built before any sample exists.  SVD identifies which Z directions a probe set constrains.  Aim 1→3 closes specific ill-conditioned directions; condition number κ tracks the budget."),

    ("All κ values are computed by SVD on the Jacobian before any sample is fabricated — experimental power is a number you can plan around, not a hope after data collection.",
     "All κ values are computed by SVD before any sample is fabricated — experimental power is a number you plan around, not a hope after data collection."),
]


# ============================================================
# S10 — Roadmap
# ============================================================
S10_REPLACEMENTS = [
    ("Year 1 builds the pipeline on synthetic data and a fiducial-grid sample platform.",
     "Year 1 builds the inference pipeline on synthetic data  +  a fiducial-grid sample platform.  Year 2 publishes the SL-PCI proof-of-concept on the Br-rich vs passivation matrix.  Year 3 closes the mechanism map across the high-Br stress regime."),
]


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation(str(SRC))
    print(f"Opened: {SRC.name}  ({len(prs.slides)} slides)\n")

    print(">>> S2 — Background: Chinese → English")
    for old, new in S2_REPLACEMENTS:
        replace_in_slide(prs.slides[1], old, new, label="S2")

    print("\n>>> S5 — Coupled-state problem: tighten")
    for old, new in S5_REPLACEMENTS:
        replace_in_slide(prs.slides[4], old, new, label="S5")

    print("\n>>> S6 — Central hypothesis: tighten")
    for old, new in S6_REPLACEMENTS:
        replace_in_slide(prs.slides[5], old, new, label="S6")

    print("\n>>> S7 — SL-PCI method")
    for old, new in S7_REPLACEMENTS:
        replace_in_slide(prs.slides[6], old, new, label="S7")

    print("\n>>> S8 — Forward models")
    for old, new in S8_REPLACEMENTS:
        replace_in_slide(prs.slides[7], old, new, label="S8")

    print("\n>>> S9 — Identifiability")
    for old, new in S9_REPLACEMENTS:
        replace_in_slide(prs.slides[8], old, new, label="S9")

    print("\n>>> S10 — Roadmap")
    for old, new in S10_REPLACEMENTS:
        replace_in_slide(prs.slides[9], old, new, label="S10")

    prs.save(str(OUT))
    print(f"\nSaved: {OUT}")


if __name__ == "__main__":
    main()
