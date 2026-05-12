"""
Build oral_exam_v3_SLPCI.pptx — each S5–S10 has its own custom layout
designed for its content, not a clone of slide 4.

Layouts:
  S5  Coupled-state problem    : top diagram (3 probes -> 1 Z) + 3 horizontal contradiction cards
  S6  Central hypothesis       : pentagon of 5 latent states (left) + 3 falsifiable sub-hypothesis cards (right)
  S7  SL-PCI method            : feature equation banner + 4-step horizontal process flow + innovation callout
  S8  Forward models           : Z feature box + 3-domain technique matrix
  S9  Identifiability budget   : 3 progressive Aim stages with arrows + Z components highlighted as resolved
  S10 Roadmap                  : Gantt timeline (M1–M36) + 3 year-summary cards + thesis claim banner
"""
from __future__ import annotations

import copy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SRC = Path("oral exam_260423_v1.0.pptx")
OUT = Path("oral exam_260423_v3_SLPCI.pptx")

# ============================================================
# Palette — academic navy/blue, matched to slides 2–4 chrome
# ============================================================
NAVY = RGBColor(0x1F, 0x3A, 0x5F)        # primary dark
BLUE = RGBColor(0x2E, 0x75, 0xB6)        # accent mid
TEAL = RGBColor(0x17, 0x6B, 0x87)        # secondary
LIGHT_BLUE = RGBColor(0xDC, 0xE6, 0xF2)  # fill light
PALE = RGBColor(0xF5, 0xF7, 0xFA)        # very light fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_900 = RGBColor(0x21, 0x25, 0x29)    # body text
GRAY_700 = RGBColor(0x49, 0x4F, 0x55)    # secondary body
GRAY_500 = RGBColor(0x86, 0x8E, 0x96)    # caption / muted
GRAY_300 = RGBColor(0xC8, 0xCE, 0xD3)    # divider
RED = RGBColor(0xC0, 0x39, 0x2B)         # accent — corner badge
GREEN = RGBColor(0x2E, 0x86, 0x55)       # resolved / success
AMBER = RGBColor(0xD9, 0x82, 0x1C)       # caution / partial

FONT = "Calibri"
MONO = "Consolas"


# ============================================================
# Primitive helpers
# ============================================================
def _set_fill(shape, color):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def _set_line(shape, color, width=Pt(0.5)):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = width


def _no_shadow(shape):
    # python-pptx doesn't expose shadow removal — drill into XML
    sp = shape._element
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    spPr = sp.find(".//a:spPr", nsmap)
    if spPr is None:
        return
    # remove existing effectLst, add empty one
    for tag in ("effectLst", "effectDag"):
        el = spPr.find(f"a:{tag}", nsmap)
        if el is not None:
            spPr.remove(el)
    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    etree.SubElement(spPr, f"{a_ns}effectLst")


def line(slide, x1, y1, x2, y2, color=NAVY, width=Pt(0.75)):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = width
    return c


def arrow(slide, x1, y1, x2, y2, color=NAVY, width=Pt(1.5), head="triangle"):
    c = line(slide, x1, y1, x2, y2, color=color, width=width)
    spPr = c.line._get_or_add_ln()
    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    # remove existing tailEnd
    for t in spPr.findall(f"{a_ns}tailEnd"):
        spPr.remove(t)
    tail = etree.SubElement(spPr, f"{a_ns}tailEnd")
    tail.set("type", head)
    tail.set("w", "med")
    tail.set("len", "med")
    return c


def rect(slide, x, y, w, h, fill=None, line=None, lw=Pt(0.5),
         shape_type=MSO_SHAPE.RECTANGLE):
    r = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(r, fill)
    _set_line(r, line, lw)
    _no_shadow(r)
    return r


def text(slide, x, y, w, h, content, size=11, color=GRAY_900, bold=False,
         italic=False, align="left", anchor="top", font=FONT, lh=1.15):
    """Add a text box; content can be str or list of (text, opts) tuples."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    if anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    lines = content if isinstance(content, list) else content.split("\n")

    for i, raw in enumerate(lines):
        # Each raw can be either a string (use defaults) or a dict
        if isinstance(raw, tuple):
            txt, opts = raw
        elif isinstance(raw, dict):
            txt = raw.get("text", "")
            opts = raw
        else:
            txt = str(raw)
            opts = {}

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(opts.get("align", align), PP_ALIGN.LEFT)
        p.line_spacing = opts.get("lh", lh)
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])

        r = p.add_run()
        r.text = txt
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.italic = opts.get("italic", italic)
        r.font.color.rgb = opts.get("color", color)
    return box


# ============================================================
# Chrome — shared frame elements (top/bottom lines, red square, page #)
# ============================================================
def chrome(slide, page_label, title_text):
    """Add chrome and slide title; returns nothing."""
    # Top divider line
    line(slide, 0, 0.6, 13.33, 0.6, color=NAVY, width=Pt(1.25))
    # Bottom divider line
    line(slide, 0, 7.2, 13.33, 7.2, color=NAVY, width=Pt(0.75))
    # Red corner badge top-right
    rect(slide, 12.95, 0.05, 0.3, 0.3, fill=RED)
    # Slide title
    text(slide, 0.1, 0.08, 12.7, 0.46, title_text,
         size=20, color=NAVY, bold=True, anchor="middle")
    # Page label bottom-right
    text(slide, 12.55, 7.25, 0.7, 0.2, page_label,
         size=10, color=GRAY_500, align="right")


# ============================================================
# S5 — Coupled-state problem
#   Layout: hero diagram (3 probes -> 1 Z) on top half,
#           3 horizontal contradiction cards on bottom half
# ============================================================
def build_s5(slide):
    chrome(slide, "0 5",
           "4. Phase segregation is a coupled-state problem — three probes, three projections, one shared latent state")

    # Subtitle banner
    text(slide, 0.2, 0.68, 13.0, 0.42,
         "Three contradictions in the literature are not measurement noise. Each technique reports a different projection of the same physical state, so their disagreements are diagnostic — not failures to be averaged away.",
         size=12, color=GRAY_700, italic=True, lh=1.2)

    # ------- Hero diagram (T=1.15 to 3.45) -------
    # 3 probe boxes at top
    probe_y = 1.30
    probe_h = 0.65
    probe_w = 3.6
    probe_gap = 0.45
    total_w = 3 * probe_w + 2 * probe_gap
    probe_x_start = (13.33 - total_w) / 2

    probes = [
        ("PL Hyperspectral", "peak, FWHM, dynamics"),
        ("KPFM / fast-KPFM", "CPD, SPV, dark relax"),
        ("XPS · TOF-SIMS · GIWAXS · RGA", "depth · chemistry · structure · gas"),
    ]
    probe_centers = []
    for i, (label, sub) in enumerate(probes):
        x = probe_x_start + i * (probe_w + probe_gap)
        rect(slide, x, probe_y, probe_w, probe_h, fill=LIGHT_BLUE, line=BLUE, lw=Pt(1.0))
        text(slide, x, probe_y + 0.04, probe_w, 0.30, label,
             size=12, color=NAVY, bold=True, align="center", anchor="top")
        text(slide, x, probe_y + 0.34, probe_w, 0.30, sub,
             size=9.5, color=GRAY_700, align="center", anchor="top", italic=True)
        probe_centers.append(x + probe_w / 2)

    # Central latent-state box
    z_y = 2.85
    z_h = 0.65
    z_w = 8.0
    z_x = (13.33 - z_w) / 2
    rect(slide, z_x, z_y, z_w, z_h, fill=NAVY, line=NAVY, lw=Pt(1.0))
    text(slide, z_x, z_y + 0.04, z_w, 0.28,
         "Shared latent state",
         size=10, color=LIGHT_BLUE, bold=True, align="center")
    text(slide, z_x, z_y + 0.30, z_w, 0.32,
         "Z  =  { x_Br/I ,  N_defect ,  φ_local ,  ε ,  C_deg }",
         size=14, color=WHITE, bold=True, align="center", font=MONO)

    # Arrows from each probe down toward central Z
    z_top_y = z_y
    for cx in probe_centers:
        arrow(slide, cx, probe_y + probe_h + 0.02, cx, z_top_y - 0.02,
              color=NAVY, width=Pt(1.5))

    # Caption under the diagram
    text(slide, 0.5, 3.62, 12.3, 0.30,
         "Forward observation model:   y_m  =  H_m ( Z ;  θ_m )  +  A_m  +  η_m       (physics · artifact channel · noise)",
         size=11, color=GRAY_700, align="center", italic=True, font=MONO)

    # ------- 3 contradiction cards (T = 4.05 to 6.9, three rows) -------
    cards = [
        dict(
            num="1",
            title="PL signal is funneling-weighted, not phase-fraction-weighted",
            body="Minority I-rich domains capture carriers from a diffusion volume V_D and dominate emission even at φ_I ≈ 10⁻³. A PL red-shift therefore tells you carriers found a low-E_g site — it does not tell you how much phase formed.",
            ref="Hoke 2015 · Ruth 2023 (ΔF_light, V_D)",
        ),
        dict(
            num="2",
            title="KPFM CPD is a 5-component convolution",
            body="V_CPD = W₀(χ_Br_surf) − V_bb(N_defect) − V_SPV(I, N_defect) − V_ion(t, μ_X) − V_trap(r, t). A single static scan cannot attribute the signal to any one component, yet field interpretations routinely read it as local electric field.",
            ref="Qu 2026 · Garrett fast-KPFM (dark relax)",
        ),
        dict(
            num="3",
            title="Ex-situ probes can drive the very state they measure",
            body="X-ray dose (XPS) and sputter (TOF-SIMS) move halides. A reported 'Br-rich surface' may be the probe's artifact, not the sample's state. The Fang 2024 vs Navid 2026 directional disagreement on Br-surface stabilization is the canonical symptom.",
            ref="Tian 2022 · Fang 2024 ⟷ Navid 2026",
        ),
    ]
    card_y_start = 4.05
    card_h = 0.92
    card_gap = 0.08
    for i, c in enumerate(cards):
        y = card_y_start + i * (card_h + card_gap)
        # Number badge — square at left
        rect(slide, 0.30, y, 0.55, card_h, fill=NAVY, line=NAVY)
        text(slide, 0.30, y, 0.55, card_h, c["num"],
             size=28, color=WHITE, bold=True, align="center", anchor="middle")
        # Card body strip
        rect(slide, 0.85, y, 12.18, card_h, fill=PALE, line=GRAY_300, lw=Pt(0.5))
        # Title row
        text(slide, 1.00, y + 0.06, 9.5, 0.28, c["title"],
             size=12, color=NAVY, bold=True)
        # Body
        text(slide, 1.00, y + 0.34, 9.5, card_h - 0.40, c["body"],
             size=10, color=GRAY_900, lh=1.15)
        # Reference (right side)
        rect(slide, 10.55, y + 0.10, 2.40, card_h - 0.20, fill=WHITE, line=BLUE, lw=Pt(0.5))
        text(slide, 10.55, y + 0.10, 2.40, 0.24, "LITERATURE",
             size=8, color=BLUE, bold=True, align="center", anchor="middle")
        text(slide, 10.55, y + 0.34, 2.40, card_h - 0.46, c["ref"],
             size=9.5, color=GRAY_900, align="center", anchor="middle", italic=True)


# ============================================================
# S6 — Central Hypothesis
#   Layout: pentagon of 5 coupled latent states (left)
#         + 3 falsifiable hypothesis cards (right)
# ============================================================
def build_s6(slide):
    chrome(slide, "0 6",
           "5. Br-rich surface is conditionally stabilizing — three falsifiable sub-hypotheses on a coupled 5-state system")

    text(slide, 0.2, 0.68, 13.0, 0.42,
         "Phase segregation is governed by five coupled latent states. Br richness is neither necessary nor sufficient — its role is set by whether it co-moves with defect, electrostatic, strain and chemical-potential changes.",
         size=12, color=GRAY_700, italic=True, lh=1.2)

    # === LEFT: pentagon diagram (L = 0.2 to 6.4) ===
    # Pentagon vertices at radius ~2.3 in a 6.2x6.0 area centered at (3.2, 4.3)
    import math
    cx, cy = 3.2, 4.3
    r = 2.0
    # 5 vertices, starting at top, going clockwise
    vertices = []
    for k in range(5):
        ang = -math.pi / 2 + k * (2 * math.pi / 5)
        vx = cx + r * math.cos(ang)
        vy = cy + r * math.sin(ang)
        vertices.append((vx, vy))

    labels = [
        ("x_Br/I", "halide distribution"),
        ("N_defect", "defect density"),
        ("φ_local", "local potential"),
        ("ε", "strain heterogeneity"),
        ("C_deg", "irreversible chemistry"),
    ]

    # Section header for left side
    text(slide, 0.30, 1.20, 6.2, 0.35, "FIVE COUPLED LATENT STATES",
         size=11, color=BLUE, bold=True, align="center")
    text(slide, 0.30, 1.50, 6.2, 0.25, "every pair coupled bidirectionally",
         size=9.5, color=GRAY_500, italic=True, align="center")

    # Draw connecting lines (complete graph K5)
    for i in range(5):
        for j in range(i + 1, 5):
            v1, v2 = vertices[i], vertices[j]
            line(slide, v1[0], v1[1], v2[0], v2[1], color=GRAY_300, width=Pt(0.5))

    # Draw vertex circles + labels
    node_w, node_h = 1.30, 0.62
    for (vx, vy), (sym, sub) in zip(vertices, labels):
        # Circle/oval
        rect(slide, vx - node_w / 2, vy - node_h / 2, node_w, node_h,
             fill=NAVY, line=NAVY, shape_type=MSO_SHAPE.OVAL)
        text(slide, vx - node_w / 2, vy - node_h / 2 - 0.02,
             node_w, 0.32, sym,
             size=12, color=WHITE, bold=True, align="center", anchor="middle",
             font=MONO)
        text(slide, vx - node_w / 2, vy - 0.02,
             node_w, 0.26, sub,
             size=8, color=LIGHT_BLUE, align="center", anchor="middle", italic=True)

    # Bottom caption for left side
    text(slide, 0.30, 6.50, 6.2, 0.30,
         "Br-rich surface acts on x_Br/I directly — but its phenotype depends on every other vertex moving with it.",
         size=10, color=GRAY_700, align="center", italic=True)

    # === RIGHT: 3 hypothesis cards (L = 6.7 to 13.05) ===
    text(slide, 6.70, 1.20, 6.30, 0.35,
         "THREE FALSIFIABLE SUB-HYPOTHESES",
         size=11, color=BLUE, bold=True, align="center")
    text(slide, 6.70, 1.50, 6.30, 0.25,
         "each with a statistical decision rule, not a truth-table",
         size=9.5, color=GRAY_500, italic=True, align="center")

    hypotheses = [
        dict(
            tag="Hα",
            name="Conditional stabilization",
            claim="Br-rich surface stabilizes only when it simultaneously reduces N_defect, smooths φ_local, and releases ε.",
            decision="Multi-variate regression slope β₁ in:\nφ = β₀ + β₁·χ_Br_surf + β₂·N_defect + β₃·ε + β₄·φ_loc + ε",
            metric="β₁ < 0,  95% CI excludes 0",
        ),
        dict(
            tag="Hβ",
            name="Electrostatic precursor",
            claim="Local CPD anomaly is a trigger only if it precedes the PL hotspot in time AND co-locates in the same ROI.",
            decision="Causal-precedence probability across\nstate-locked S0→S5 sequence in same ROI",
            metric="P[t_CPD < t_PL | ROI] > 0.5 + δ",
        ),
        dict(
            tag="Hγ",
            name="Reversibility boundary",
            claim="Reversible LHS and irreversible degradation share PL signatures — separable only by joint posterior across optical + chemical + gas channels.",
            decision="Bayes-factor K between\nM_reversible vs M_irreversible models",
            metric="|log₁₀ K|  >  1   ⟹  decisive",
        ),
    ]
    card_y_start = 1.95
    card_h = 1.65
    card_gap = 0.10
    for i, h in enumerate(hypotheses):
        y = card_y_start + i * (card_h + card_gap)
        # outer card
        rect(slide, 6.70, y, 6.30, card_h, fill=PALE, line=GRAY_300, lw=Pt(0.5))
        # tag column on left
        rect(slide, 6.70, y, 0.85, card_h, fill=NAVY, line=NAVY)
        text(slide, 6.70, y + 0.05, 0.85, 0.55, h["tag"],
             size=22, color=WHITE, bold=True, align="center", anchor="middle")
        text(slide, 6.70, y + 0.62, 0.85, 0.30, h["name"].upper().split(" ")[0],
             size=7.5, color=LIGHT_BLUE, bold=True, align="center", anchor="middle")
        # name header
        text(slide, 7.65, y + 0.05, 5.30, 0.30, h["name"],
             size=12.5, color=NAVY, bold=True)
        # claim
        text(slide, 7.65, y + 0.34, 5.30, 0.62, h["claim"],
             size=9.5, color=GRAY_900, lh=1.15)
        # decision line
        text(slide, 7.65, y + 0.99, 5.30, 0.38, h["decision"],
             size=8.5, color=GRAY_700, italic=True, lh=1.1)
        # metric strip
        rect(slide, 7.65, y + 1.34, 5.30, 0.26, fill=LIGHT_BLUE, line=None)
        text(slide, 7.65, y + 1.34, 5.30, 0.26, "DECISION:  " + h["metric"],
             size=9.5, color=NAVY, bold=True, align="center", anchor="middle", font=MONO)


# ============================================================
# S7 — SL-PCI method: 4-step process flow
#   Layout: feature equation banner on top, 4 numbered cards
#           connected by arrows in middle, innovation callout bottom
# ============================================================
def build_s7(slide):
    chrome(slide, "0 7",
           "6. SL-PCI: four steps that turn correlative observation into Bayesian inference")

    text(slide, 0.2, 0.68, 13.0, 0.42,
         "State-Locked Physics-Informed Correlative Inference — write a defined state, freeze it for endpoint probes, infer a shared latent state with explicit artifact channels.",
         size=12, color=GRAY_700, italic=True, lh=1.2)

    # === Feature equation banner (T=1.15 to 2.00) ===
    rect(slide, 0.5, 1.15, 12.33, 0.85, fill=NAVY, line=NAVY)
    text(slide, 0.5, 1.18, 12.33, 0.30, "GENERATIVE OBSERVATION MODEL",
         size=10, color=LIGHT_BLUE, bold=True, align="center", anchor="middle")
    text(slide, 0.5, 1.45, 12.33, 0.40,
         "y_m  =  H_m (  Z  ;  θ_m )  +  A_m  +  η_m",
         size=22, color=WHITE, bold=True, align="center", anchor="middle", font=MONO)
    text(slide, 0.5, 1.80, 12.33, 0.22,
         "observation   =   physics forward model   +   artifact channel   +   noise",
         size=10, color=LIGHT_BLUE, italic=True, align="center", anchor="middle")

    # === 4-step process flow (T = 2.40 to 5.50) ===
    steps = [
        dict(
            n="1",
            head="STATE WRITING",
            sub="define · timestamp",
            body="Operando PL + KPFM in a fiducial ROI define six states S0–S5 from peak shift, FWHM and CPD dynamics. Each measurement carries a full history H = {I, t, T, atm, delay, dose}.",
            tech="PL hyperspectral · fast-KPFM",
            tscale="ms — min",
        ),
        dict(
            n="2",
            head="STATE LOCKING",
            sub="freeze · transport",
            body="Low-T quench + dark transfer + inert atmosphere keep the written state frozen during transit to destructive probes. Twin-sample dose array provides the backup when cryo-transfer is unavailable.",
            tech="cryo-stage · twin-sample",
            tscale="seconds (quench)",
        ),
        dict(
            n="3",
            head="ENDPOINT READOUT",
            sub="chemistry · structure",
            body="TOF-SIMS / AR-XPS / GIWAXS / Raman / RGA read the locked state. Each instrument's artifact channel A_m is calibrated by dose / sputter / exposure series — not assumed to be zero.",
            tech="SIMS · XPS · GIWAXS · RGA",
            tscale="static endpoint",
        ),
        dict(
            n="4",
            head="LATENT INFERENCE",
            sub="posterior · uncertainty",
            body="Bayesian inference returns p(Z, θ | y) with covariance. Jacobian SVD identifies ill-conditioned directions before any sample is made. Output: mechanism posteriors and Bayes factors between hypotheses.",
            tech="HMC / VI · Jacobian SVD",
            tscale="offline",
        ),
    ]
    step_y = 2.55
    step_h = 2.85
    step_w = 2.80
    step_gap = 0.30
    total = 4 * step_w + 3 * step_gap
    x0 = (13.33 - total) / 2

    for i, s in enumerate(steps):
        x = x0 + i * (step_w + step_gap)
        # Card outline
        rect(slide, x, step_y, step_w, step_h, fill=PALE, line=BLUE, lw=Pt(0.75))
        # Header band
        rect(slide, x, step_y, step_w, 0.78, fill=NAVY, line=NAVY)
        # Number circle on header
        rect(slide, x + 0.10, step_y + 0.12, 0.55, 0.55,
             fill=BLUE, line=BLUE, shape_type=MSO_SHAPE.OVAL)
        text(slide, x + 0.10, step_y + 0.12, 0.55, 0.55, s["n"],
             size=18, color=WHITE, bold=True, align="center", anchor="middle")
        # Head & sub
        text(slide, x + 0.70, step_y + 0.10, step_w - 0.80, 0.28, s["head"],
             size=11, color=WHITE, bold=True)
        text(slide, x + 0.70, step_y + 0.40, step_w - 0.80, 0.24, s["sub"],
             size=8.5, color=LIGHT_BLUE, italic=True)

        # Body
        text(slide, x + 0.13, step_y + 0.92, step_w - 0.26, 1.40, s["body"],
             size=9.5, color=GRAY_900, lh=1.20)

        # Bottom strip — technique + timescale
        rect(slide, x, step_y + step_h - 0.55, step_w, 0.55,
             fill=LIGHT_BLUE, line=None)
        text(slide, x + 0.13, step_y + step_h - 0.52, step_w - 0.26, 0.22,
             "TECHNIQUE:  " + s["tech"], size=8.5, color=NAVY, bold=True)
        text(slide, x + 0.13, step_y + step_h - 0.30, step_w - 0.26, 0.22,
             "TIME WINDOW:  " + s["tscale"], size=8.5, color=NAVY)

        # Arrow to next step
        if i < 3:
            ay = step_y + step_h / 2
            arrow(slide, x + step_w + 0.03, ay,
                  x + step_w + step_gap - 0.03, ay,
                  color=NAVY, width=Pt(2.0))

    # === Innovation callout (T = 5.65 to 6.95) ===
    rect(slide, 0.5, 5.65, 12.33, 1.30, fill=LIGHT_BLUE, line=BLUE, lw=Pt(0.75))
    text(slide, 0.7, 5.72, 11.93, 0.32, "WHY THIS IS NEW",
         size=11, color=NAVY, bold=True)
    text(slide, 0.7, 6.04, 11.93, 0.85,
         [("•  No technique is asked to measure Z directly — each contributes a noisy, biased projection (H_m + A_m).", {}),
          ("•  The artifact channel A_m is explicit in the likelihood, not denied — beam damage and sputter mixing are first-class citizens.", {}),
          ("•  Jacobian SVD on the chosen probe set predicts identifiability before any sample is fabricated — experimental power is a quantitative number, not a hope.", {})],
         size=10.5, color=GRAY_900, lh=1.30)


# ============================================================
# S8 — Forward Models matrix
# ============================================================
def build_s8(slide):
    chrome(slide, "0 8",
           "7. Forward models link nine techniques to one shared latent state Z")

    text(slide, 0.2, 0.68, 13.0, 0.42,
         "Listing H_m(Z; θ_m) explicitly makes the inference reproducible. Each technique's primary observable, governing equation and resolved Z components are laid out before any sample is fabricated.",
         size=12, color=GRAY_700, italic=True, lh=1.2)

    # Z definition feature bar
    rect(slide, 0.50, 1.15, 12.33, 0.55, fill=NAVY, line=NAVY)
    text(slide, 0.50, 1.15, 5.0, 0.55, "  LATENT STATE",
         size=11, color=LIGHT_BLUE, bold=True, anchor="middle")
    text(slide, 4.50, 1.15, 8.33, 0.55,
         "Z = { x_Br/I ,  N_defect ,  φ_local ,  ε ,  C_deg }",
         size=15, color=WHITE, bold=True, align="center", anchor="middle", font=MONO)

    # === 3-domain table ===
    # Columns: Technique | Observable | Forward model (excerpt) | Resolves
    # Domains: Optical (3 rows), Electrostatic (1 row), Chemical/Structural (3 rows)
    domains = [
        dict(
            tag="OPTICAL  PROJECTIONS",
            color=BLUE,
            rows=[
                ("PL hyperspectral", "I_PL(E, r, t)",
                 "I_PL = Σ_p  w_p · η_p(N_defect) · L_p(E_g^p, σ) · R(α,d)\n         w_p ∝ V_D⁻¹ · exp(−ΔE_g/kT)    (Ruth 2023 funneling)",
                 "x_Br/I , φ_p , N_defect"),
                ("PLQY  +  TRPL", "η_PLQY ,  τ_eff",
                 "η⁻¹ − 1  =  1/(k_rad · τ_SRH(N_b))  +  S_surf(N_s)/(k_rad · d)\n         τ_eff⁻¹  =  τ_SRH⁻¹  +  τ_rad⁻¹  +  τ_surf⁻¹",
                 "N_defect (bulk / surf)"),
            ],
        ),
        dict(
            tag="ELECTROSTATIC  PROJECTIONS",
            color=TEAL,
            rows=[
                ("KPFM  +  fast-KPFM", "V_CPD (r, t ;  I,  f )",
                 "V_CPD  =  W₀(χ_Br^surf)  −  V_bb(N_defect)  −  V_SPV(I, N_defect)\n               −  V_ion(t, μ_X)  −  V_trap(r, t)        (separable by τ:  ms / s / min)",
                 "χ_Br^surf , N_defect , φ_local"),
            ],
        ),
        dict(
            tag="CHEMICAL  +  STRUCTURAL  PROJECTIONS",
            color=AMBER,
            rows=[
                ("AR-XPS", "I_Br(θ) / I_I(θ)",
                 "I_Br(θ) / I_I(θ)  =  ∫ c_Br(z) · exp(−z / λ cos θ) dz  /  ∫ c_I(...)\n         dose-series  →  zero-dose extrapolation removes beam-driven A_m",
                 "x_Br/I (z) , C_deg"),
                ("GIWAXS  +  Raman", "Δq , FWHM , ω, Γ",
                 "q_hkl(x_Br) = q₀(1 − x_Br · Δa/a)        (Vegard)\n         FWHM  ≈  ε · tan θ  +  K λ / L        (Williamson–Hall)",
                 "x_Br/I , ε , L_crystallite"),
                ("RGA  +  sealed-vs-open", "P_I₂(t)",
                 "dN_I₂/dt  =  R_loss(c_I , T , hν) · V_film\n         atmosphere-control isolates irreversible loss from reversible drift",
                 "C_deg ,  μ_X"),
            ],
        ),
    ]

    # Table dimensions
    table_x = 0.50
    table_y_start = 1.95
    table_w = 12.33
    col_w = [2.50, 1.95, 5.40, 2.48]  # Technique / Observable / Forward model / Resolves

    # Header row
    header_h = 0.32
    rect(slide, table_x, table_y_start, table_w, header_h,
         fill=GRAY_900, line=GRAY_900)
    headers = ["TECHNIQUE", "OBSERVABLE", "FORWARD MODEL  H_m(Z ;  θ_m)", "RESOLVES IN Z"]
    cx = table_x
    for i, hd in enumerate(headers):
        text(slide, cx, table_y_start, col_w[i], header_h, hd,
             size=9, color=LIGHT_BLUE, bold=True, align="center", anchor="middle")
        cx += col_w[i]

    # Body rows
    cur_y = table_y_start + header_h
    for d in domains:
        # Domain tag row
        rect(slide, table_x, cur_y, table_w, 0.26,
             fill=d["color"], line=d["color"])
        text(slide, table_x + 0.15, cur_y, table_w - 0.30, 0.26, d["tag"],
             size=10, color=WHITE, bold=True, anchor="middle")
        cur_y += 0.26
        for row in d["rows"]:
            tech, obs, model, resolves = row
            # row height: depends on lines in model
            n_lines_model = model.count("\n") + 1
            row_h = max(0.62, 0.30 * n_lines_model + 0.14)

            # Row background alternating
            rect(slide, table_x, cur_y, table_w, row_h,
                 fill=WHITE, line=GRAY_300, lw=Pt(0.4))
            # column dividers
            cx = table_x
            for cw in col_w[:-1]:
                cx += cw
                line(slide, cx, cur_y, cx, cur_y + row_h, color=GRAY_300, width=Pt(0.4))

            cx = table_x
            # Technique
            text(slide, cx + 0.10, cur_y + 0.05, col_w[0] - 0.15, row_h - 0.10,
                 tech, size=10, color=NAVY, bold=True, anchor="middle")
            cx += col_w[0]
            # Observable
            text(slide, cx + 0.10, cur_y + 0.05, col_w[1] - 0.15, row_h - 0.10,
                 obs, size=9.5, color=GRAY_900, anchor="middle", font=MONO)
            cx += col_w[1]
            # Forward model (multi-line, monospace)
            text(slide, cx + 0.10, cur_y + 0.05, col_w[2] - 0.15, row_h - 0.10,
                 model, size=9, color=GRAY_900, anchor="middle", font=MONO, lh=1.15)
            cx += col_w[2]
            # Resolves
            text(slide, cx + 0.10, cur_y + 0.05, col_w[3] - 0.15, row_h - 0.10,
                 resolves, size=9.5, color=NAVY, bold=True, anchor="middle", italic=False)
            cur_y += row_h

    # Footer
    text(slide, 0.50, 6.90, 12.33, 0.25,
         "Every probe is over-determined when combined — no single instrument resolves Z, but their joint likelihood does.",
         size=10, color=GRAY_700, align="center", italic=True)


# ============================================================
# S9 — Identifiability budget
#   Layout: 3 Aim stages connected by arrows, each showing
#   progressively more resolved Z components and κ improving
# ============================================================
def build_s9(slide):
    chrome(slide, "0 9",
           "8. Each Aim progressively closes the unresolved Z directions — κ tracks the budget")

    text(slide, 0.2, 0.68, 13.0, 0.42,
         "Jacobian J_mj = ∂H_m / ∂Z_j is built from the forward models before any sample exists. Its SVD reveals which mechanism directions a probe set constrains. Aim 1→3 adds probes that close specific ill-conditioned directions; condition number κ tracks progress.",
         size=11.5, color=GRAY_700, italic=True, lh=1.2)

    # === Three Aim stages, left-to-right with arrows ===
    Z_components = ["x_Br/I", "N_defect_b", "N_defect_s", "φ_local", "ε", "L_xtal", "C_deg"]

    # Aim 1 resolves: x_Br_bulk (~partial), so coarse x_Br/I only
    # Aim 2 adds: fast-KPFM 3D + TOF-SIMS + AR-XPS  -> resolves φ_local, N_defect_surf, x_Br(z)
    # Aim 3 adds: GIWAXS+Raman + RGA  -> resolves epsilon, L, C_deg

    aims = [
        dict(
            tag="AIM 1",
            year="Year 1",
            sub="Foundation",
            probes=[
                "PL hyperspectral",
                "UV-vis",
                "Static KPFM",
                "UPS (W₀ anchor)",
                "M0 synthetic-data pipeline",
            ],
            resolved={"x_Br/I": "partial", "N_defect_b": "partial",
                      "N_defect_s": "open", "φ_local": "open",
                      "ε": "open", "L_xtal": "open", "C_deg": "open"},
            kappa="κ  =  baseline",
            kappa_note="N_b vs N_s degenerate;\nC_deg unobservable",
        ),
        dict(
            tag="AIM 2",
            year="Year 2",
            sub="Core proof",
            probes=[
                "+ fast-KPFM 3D scan",
                "+ TOF-SIMS depth",
                "+ AR-XPS",
                "+ TRPL thickness gradient",
                "all in same fiducial ROI",
            ],
            resolved={"x_Br/I": "ok", "N_defect_b": "ok",
                      "N_defect_s": "ok", "φ_local": "ok",
                      "ε": "open", "L_xtal": "open", "C_deg": "open"},
            kappa="κ  ≤  10²",
            kappa_note="ε vs L still confounded;\nC_deg implicit only",
        ),
        dict(
            tag="AIM 3",
            year="Year 3",
            sub="Mechanism map",
            probes=[
                "+ GIWAXS  (Δq , FWHM)",
                "+ Raman   (linewidth)",
                "+ RGA / gas trapping",
                "+ sealed-vs-open aging",
                "high-Br stress regime",
            ],
            resolved={"x_Br/I": "ok", "N_defect_b": "ok",
                      "N_defect_s": "ok", "φ_local": "ok",
                      "ε": "ok", "L_xtal": "ok", "C_deg": "ok"},
            kappa="κ  ≈  10",
            kappa_note="all Z resolved;\nHα / Hβ / Hγ testable",
        ),
    ]

    stage_y = 1.32
    stage_h = 5.20
    stage_w = 4.00
    stage_gap = 0.35
    total = 3 * stage_w + 2 * stage_gap
    x0 = (13.33 - total) / 2

    status_color = {"ok": GREEN, "partial": AMBER, "open": GRAY_300}
    status_text_color = {"ok": WHITE, "partial": WHITE, "open": GRAY_500}
    status_glyph = {"ok": "✓", "partial": "◐", "open": "○"}

    for i, a in enumerate(aims):
        x = x0 + i * (stage_w + stage_gap)
        # Outer card
        rect(slide, x, stage_y, stage_w, stage_h, fill=PALE, line=NAVY, lw=Pt(0.8))
        # Header band
        rect(slide, x, stage_y, stage_w, 0.78, fill=NAVY, line=NAVY)
        text(slide, x + 0.15, stage_y + 0.05, stage_w - 0.30, 0.32, a["tag"],
             size=14, color=WHITE, bold=True)
        text(slide, x + 0.15, stage_y + 0.36, stage_w - 0.30, 0.22,
             a["year"] + "   ·   " + a["sub"],
             size=10, color=LIGHT_BLUE, italic=True)

        # Section: probes added
        text(slide, x + 0.15, stage_y + 0.90, stage_w - 0.30, 0.25,
             "PROBE SET",
             size=9, color=BLUE, bold=True)
        # bullets
        for k, p in enumerate(a["probes"]):
            text(slide, x + 0.20, stage_y + 1.15 + k * 0.22, stage_w - 0.40, 0.22,
                 "•  " + p, size=9.5, color=GRAY_900)

        # Section: Z resolution checklist
        check_y = stage_y + 2.45
        text(slide, x + 0.15, check_y, stage_w - 0.30, 0.25,
             "Z  DIMENSIONS",
             size=9, color=BLUE, bold=True)
        for k, z_name in enumerate(Z_components):
            yy = check_y + 0.27 + k * 0.27
            st = a["resolved"][z_name]
            # Status pill
            rect(slide, x + 0.20, yy, 0.32, 0.22,
                 fill=status_color[st], line=status_color[st])
            text(slide, x + 0.20, yy, 0.32, 0.22, status_glyph[st],
                 size=10, color=status_text_color[st], bold=True,
                 align="center", anchor="middle")
            # Name
            text(slide, x + 0.58, yy, stage_w - 0.65, 0.22, z_name,
                 size=9.5,
                 color=GRAY_900 if st != "open" else GRAY_500,
                 bold=(st == "ok"), anchor="middle", font=MONO)

        # κ metric — bottom strip
        kappa_y = stage_y + stage_h - 0.80
        rect(slide, x, kappa_y, stage_w, 0.80, fill=BLUE, line=BLUE)
        text(slide, x + 0.15, kappa_y + 0.05, stage_w - 0.30, 0.34, a["kappa"],
             size=14, color=WHITE, bold=True, align="center", anchor="middle", font=MONO)
        text(slide, x + 0.15, kappa_y + 0.40, stage_w - 0.30, 0.36, a["kappa_note"],
             size=8.5, color=LIGHT_BLUE, italic=True, align="center", anchor="middle", lh=1.1)

        # Arrow to next stage
        if i < 2:
            ay = stage_y + stage_h / 2
            arrow(slide, x + stage_w + 0.05, ay,
                  x + stage_w + stage_gap - 0.05, ay,
                  color=NAVY, width=Pt(2.5))
            # "+" tag on arrow
            text(slide, x + stage_w + 0.02, ay - 0.32, stage_gap, 0.25,
                 "+ probes", size=8, color=NAVY, italic=True, bold=True, align="center")

    # Footer — SVD pre-experiment design callout
    text(slide, 0.5, 6.65, 12.33, 0.40,
         "All κ values are computed by SVD on the Jacobian before any sample is fabricated — experimental power is a number you can plan around, not a hope after data collection.",
         size=10.5, color=GRAY_700, align="center", italic=True)


# ============================================================
# S10 — Roadmap (Gantt timeline + thesis claim)
# ============================================================
def build_s10(slide):
    chrome(slide, "1 0",
           "9. Three-year roadmap — two manuscripts and one thesis, all anchored on falsifiable claims with reported uncertainty")

    text(slide, 0.2, 0.68, 13.0, 0.36,
         "Year 1 builds the pipeline on synthetic data and a fiducial-grid sample platform. Year 2 publishes the SL-PCI proof-of-concept on the Br-rich vs passivation matrix. Year 3 closes the mechanism map across the high-Br stress regime and writes the thesis.",
         size=11.5, color=GRAY_700, italic=True, lh=1.2)

    # === Gantt timeline (T = 1.20 to 4.80) ===
    timeline_x0 = 1.40
    timeline_y0 = 1.30
    timeline_w = 11.40
    timeline_h = 3.55
    month_count = 36
    per_month = timeline_w / month_count

    # Header row: months
    text(slide, 0.20, timeline_y0, 1.10, 0.30, "Quarter",
         size=9, color=GRAY_500, bold=True, anchor="middle")
    # Quarter ticks (every 3 months)
    for q in range(13):
        m = q * 3
        x = timeline_x0 + m * per_month
        line(slide, x, timeline_y0, x, timeline_y0 + timeline_h,
             color=GRAY_300, width=Pt(0.4))
        # Label
        if q == 0:
            label = "M1"
        else:
            label = f"M{m}"
        text(slide, x - 0.35, timeline_y0, 0.7, 0.30, label,
             size=8.5, color=GRAY_500, align="center", anchor="middle")

    # Year separator lines (M12, M24, M36)
    for ym in (12, 24):
        x = timeline_x0 + ym * per_month
        line(slide, x, timeline_y0 + 0.30, x, timeline_y0 + timeline_h,
             color=NAVY, width=Pt(1.5))

    # Year labels above swimlanes
    text(slide, timeline_x0, timeline_y0 + 0.32,
         12 * per_month, 0.26, "YEAR 1 — FOUNDATION",
         size=10, color=NAVY, bold=True, align="center")
    text(slide, timeline_x0 + 12 * per_month, timeline_y0 + 0.32,
         12 * per_month, 0.26, "YEAR 2 — CORE PROOF",
         size=10, color=NAVY, bold=True, align="center")
    text(slide, timeline_x0 + 24 * per_month, timeline_y0 + 0.32,
         12 * per_month, 0.26, "YEAR 3 — MECHANISM MAP",
         size=10, color=NAVY, bold=True, align="center")

    # Swimlane rows
    swimlanes = [
        dict(label="M0 · Inference pipeline\n(synthetic data)", row_y=0.70,
             bars=[(0, 4, BLUE)]),
        dict(label="Sample platform\n(fiducial-grid WBG films)", row_y=1.10,
             bars=[(1, 12, NAVY), (12, 24, NAVY)]),
        dict(label="State writing & locking\n(PL · KPFM · cryo)", row_y=1.50,
             bars=[(4, 12, TEAL)]),
        dict(label="Same-ROI multimodal\n(KPFM · TOF-SIMS · XPS)", row_y=1.90,
             bars=[(13, 24, TEAL)]),
        dict(label="Structural / chemical\n(GIWAXS · Raman · RGA)", row_y=2.30,
             bars=[(25, 33, AMBER)]),
        dict(label="Manuscript writing\n+ thesis", row_y=2.70,
             bars=[(22, 24, GREEN), (30, 36, GREEN)]),
    ]
    bar_h = 0.32
    lane_label_x = 0.20
    lane_label_w = 1.15
    for s in swimlanes:
        y = timeline_y0 + s["row_y"]
        # Lane label
        text(slide, lane_label_x, y - 0.04, lane_label_w, bar_h + 0.08, s["label"],
             size=8, color=GRAY_700, anchor="middle", lh=1.1)
        # Background lane band
        rect(slide, timeline_x0, y + 0.02, timeline_w, bar_h - 0.04,
             fill=PALE, line=None)
        # Bars
        for (m_start, m_end, color) in s["bars"]:
            bx = timeline_x0 + m_start * per_month
            bw = (m_end - m_start) * per_month
            rect(slide, bx, y, bw, bar_h, fill=color, line=color, lw=Pt(0.3))

    # Manuscript milestone diamonds (M24, M36)
    diamond_y = timeline_y0 + 3.20
    for m_label, m_pos, label in [
        ("Manuscript 1", 24, "SL-PCI proof"),
        ("Manuscript 2", 36, "Mechanism regime"),
    ]:
        x = timeline_x0 + m_pos * per_month
        rect(slide, x - 0.18, diamond_y, 0.36, 0.30, fill=RED, line=RED,
             shape_type=MSO_SHAPE.DIAMOND)
        text(slide, x - 1.10, diamond_y + 0.35, 2.20, 0.20, m_label,
             size=9, color=RED, bold=True, align="center")
        text(slide, x - 1.10, diamond_y + 0.54, 2.20, 0.18, label,
             size=8.5, color=GRAY_700, align="center", italic=True)

    # === Year summary cards (T = 5.10 to 6.55) ===
    year_y = 5.10
    year_h = 1.45
    year_w = 4.05
    year_gap = 0.20
    year_x0 = (13.33 - 3 * year_w - 2 * year_gap) / 2
    summaries = [
        dict(tag="YEAR 1", color=BLUE,
             items=[
                 "M0 synthetic pipeline validates inference recovery",
                 "fiducial-grid WBG films & baseline characterization SOP",
                 "S0–S5 state writing protocol; locking-drift quantified",
             ],
             metric="Δ_drift  <  Δ_mechanism"),
        dict(tag="YEAR 2", color=TEAL,
             items=[
                 "same-ROI PL → KPFM → TOF-SIMS on 2 × 3 matrix",
                 "leave-one-modality-out cross-validation",
                 "Manuscript 1 published — Hα tested with reported β₁",
             ],
             metric="κ  ≤  10²    ·    Hα tested"),
        dict(tag="YEAR 3", color=AMBER,
             items=[
                 "GIWAXS + Raman + RGA close ε / L / C_deg",
                 "mechanism regime diagram for E_g = 1.66–1.78 eV",
                 "Manuscript 2 + thesis with Bayes-factor support",
             ],
             metric="all Z resolved   ·   log₁₀K reported"),
    ]
    for i, s in enumerate(summaries):
        x = year_x0 + i * (year_w + year_gap)
        rect(slide, x, year_y, year_w, year_h, fill=PALE, line=s["color"], lw=Pt(0.75))
        rect(slide, x, year_y, year_w, 0.30, fill=s["color"], line=s["color"])
        text(slide, x + 0.15, year_y, year_w - 0.30, 0.30, s["tag"],
             size=10, color=WHITE, bold=True, anchor="middle")
        for k, item in enumerate(s["items"]):
            text(slide, x + 0.18, year_y + 0.38 + k * 0.22, year_w - 0.30, 0.22,
                 "•  " + item, size=9, color=GRAY_900, lh=1.10)
        # Metric strip
        rect(slide, x + 0.10, year_y + year_h - 0.30, year_w - 0.20, 0.24,
             fill=LIGHT_BLUE, line=None)
        text(slide, x + 0.10, year_y + year_h - 0.30, year_w - 0.20, 0.24,
             s["metric"], size=9, color=NAVY, bold=True,
             align="center", anchor="middle", font=MONO)

    # === Defensible thesis claim banner (T = 6.70 to 7.10) ===
    rect(slide, 0.50, 6.70, 12.33, 0.40, fill=NAVY, line=NAVY)
    text(slide, 0.50, 6.70, 12.33, 0.40,
         '"Br-rich surface is neither inherently good nor bad — its effect on phase segregation is decided by the coupled latent state it creates."',
         size=11, color=WHITE, bold=True, italic=True,
         align="center", anchor="middle")


# ============================================================
# Slide management — remove old S5–S10, insert 6 fresh ones in their place
# ============================================================
def remove_slide(prs, idx):
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    target = slides_list[idx]
    rId = target.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
    prs.part.drop_rel(rId)
    xml_slides.remove(target)


def add_blank_slide(prs):
    """Add a fresh blank slide using the same layout as slide 4 (DEFAULT)."""
    # Use the layout slide 4 uses
    layout = prs.slides[3].slide_layout
    new_slide = prs.slides.add_slide(layout)
    # Remove all auto-inserted placeholders so we have a true blank canvas
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    return new_slide


def main():
    prs = Presentation(str(SRC))
    print(f"Initial: {len(prs.slides)} slides")

    # Step 1: remove S5–S10 (indices 4–9)
    for idx in [9, 8, 7, 6, 5, 4]:
        if idx < len(prs.slides):
            remove_slide(prs, idx)
    print(f"After removal: {len(prs.slides)}")

    # Step 2: add 6 blank slides at the END
    new_slides = []
    for i in range(6):
        new_slides.append(add_blank_slide(prs))
    print(f"After append: {len(prs.slides)}")

    # Step 3: reorder — move the 6 new slides to positions 5–10
    xml_slides = prs.slides._sldIdLst
    sld_list = list(xml_slides)
    head = sld_list[:4]
    new6 = sld_list[-6:]
    middle = sld_list[4:-6]
    new_order = head + new6 + middle
    for el in sld_list:
        xml_slides.remove(el)
    for el in new_order:
        xml_slides.append(el)
    print(f"Final: {len(prs.slides)} slides, reordered")

    # Step 4: build each slide
    builders = [build_s5, build_s6, build_s7, build_s8, build_s9, build_s10]
    for i, builder in enumerate(builders):
        slide_idx = 4 + i
        print(f"  Building S{slide_idx + 1} via {builder.__name__}")
        builder(prs.slides[slide_idx])

    prs.save(str(OUT))
    print(f"\nSaved: {OUT}")


if __name__ == "__main__":
    main()
