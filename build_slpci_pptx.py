"""Build oral_exam SL-PCI v2 pptx by cloning slide 4 layout for slides 5-10."""
from __future__ import annotations

import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from lxml import etree

SRC = Path("oral exam_260423_v1.0.pptx")
OUT = Path("oral exam_260423_v2_SLPCI.pptx")


# ---------------------------------------------------------------------------
# Slide content (S5-S10): action title, intro, then 3 columns
# Each column: num, sub, fig_label, fig_ratio, fig_caption, problem_label,
# problem_body, directions_label, directions_body, metric_value, metric_caption
# Page-number text for bottom-right is set per slide.
# ---------------------------------------------------------------------------
SLIDES = [
    # S5
    dict(
        title="4. Phase Segregation is a Coupled-State Problem, Not a Single-Mechanism Failure",
        intro=(
            "Recent wide-bandgap perovskite literature exposes three unresolved tensions in how phase "
            "segregation is characterized. They are not noise: each technique probes a different physical "
            "projection of one shared latent state, and treating their readings as if they measured the "
            "same quantity is exactly why mechanism attributions remain contradictory."
        ),
        page="0 5",
        cols=[
            dict(
                num="I.", sub="PL Funneling",
                fig_label="FIG. 1", fig_ratio="ref",
                fig_caption="Hoke 2015 PL evolution; Ruth 2023 funneling free-energy schematic.",
                problem=(
                    "Minority I-rich domains dominate PL emission even at phi_I ~ 10^-3 because carriers "
                    "thermalize into the lowest-Eg sites within a diffusion volume V_D (Hoke 2015; Ruth 2023). "
                    "PL red-shift is therefore not a phase-fraction estimator."
                ),
                directions=(
                    "Funneling-corrected inversion: combine PL hyperspectral + TRPL diffusion length + GIWAXS structural fraction.\n"
                    "Use Ruth 2023 model E_PL(x)=E0(1-x)+E1*x-b*x(1-x) and weight w_p propto V_D^-1 exp(-dEg/kT) to recover phi_I.\n"
                    "Cross-check with absorption tail (Urbach E_U) to discriminate trap emission from real low-Eg phase."
                ),
                metric="phi_I uncertainty < 10%",
                metric_caption="FUNNELING-CORRECTED PHASE FRACTION TARGET",
            ),
            dict(
                num="II.", sub="KPFM Multi-Component",
                fig_label="FIG. 2", fig_ratio="ref",
                fig_caption="Qu 2026 CPD decomposition; Garrett fast-KPFM dark relaxation.",
                problem=(
                    "V_CPD = W0 - V_bb - V_SPV - V_ion - V_trap. A single static KPFM scan cannot isolate any "
                    "one component, yet field interpretations routinely equate excess CPD with local electric "
                    "field (Qu 2026)."
                ),
                directions=(
                    "Time-frequency-intensity 3D KPFM scan: light step (V_SPV ~ ms), scan-rate sweep (V_trap ~ s), dark relaxation (V_ion ~ min).\n"
                    "Anchor W0 with static UPS work-function measurement.\n"
                    "Five-component decomposition feeds the latent inference as a structured electrostatic prior."
                ),
                metric="5 components decoupled",
                metric_caption="V_bb / V_SPV / V_trap / V_ion / W0",
            ),
            dict(
                num="III.", sub="Ex-situ Drift",
                fig_label="FIG. 3", fig_ratio="ref",
                fig_caption="Tian 2022 beam-damage review; Fang 2024 vs Navid 2026 Br-surface contradiction.",
                problem=(
                    "X-ray dose (XPS) and sputter (TOF-SIMS) can themselves drive halide redistribution. A reported "
                    "'Br-rich surface' may be the probe's artifact, not the sample's state (Tian 2022; the Fang "
                    "2024 vs Navid 2026 directional disagreement is the canonical symptom)."
                ),
                directions=(
                    "State locking: low-T quench + dark transfer + inert atmosphere keeps state during transit.\n"
                    "Dose-series extrapolation to zero dose isolates intrinsic from probe-driven Br/I redistribution.\n"
                    "Artifact channel A_m enters the likelihood explicitly, so reported posteriors are debiased."
                ),
                metric="drift < signal",
                metric_caption="STATE-LOCK FIDELITY TARGET",
            ),
        ],
    ),
    # S6
    dict(
        title="5. Central Hypothesis: Br-rich Surface is Conditionally Stabilizing",
        intro=(
            "Wide-bandgap mixed-halide perovskite phase segregation is governed by five coupled latent states: "
            "halide distribution x_Br/I, defect density N_defect, local potential phi_local, strain "
            "heterogeneity epsilon, and irreversible chemistry C_deg. Br-rich surfaces are neither universally "
            "helpful nor harmful: their role is set by how they reshape the other four latent dimensions. "
            "Three sub-hypotheses make this testable."
        ),
        page="0 6",
        cols=[
            dict(
                num="H_a.", sub="Conditional Stabilization",
                fig_label="FIG. 1", fig_ratio="ref",
                fig_caption="Coupled-mechanism diagram: Br-rich <-> N_defect <-> phi <-> epsilon <-> mu_X.",
                problem=(
                    "Br-rich surface stabilizes the film only when it simultaneously reduces N_defect, smooths "
                    "phi_local fluctuations, and releases local strain epsilon. Br richness in isolation is "
                    "neither necessary nor sufficient."
                ),
                directions=(
                    "Multi-variate regression after Aim 2-3: phi = b0 + b1*x_Br_surf + b2*N_defect + b3*epsilon + b4*phi_loc + residual.\n"
                    "Decision: b1 < 0 with 95% CI excluding zero supports H_alpha; otherwise treatment is spurious.\n"
                    "Probes: AR-XPS + TOF-SIMS for x_Br_surf; PLQY + TRPL thickness gradient for N_defect."
                ),
                metric="d_phi / d_chi_Br < 0",
                metric_caption="REGRESSION SLOPE WITH 95% CI EXCLUDING ZERO",
            ),
            dict(
                num="H_b.", sub="Electrostatic Precursor",
                fig_label="FIG. 2", fig_ratio="ref",
                fig_caption="State timeline: pre-onset CPD anomaly preceding PL hotspot (Ma 2024 dark seg).",
                problem=(
                    "Local electric field is a candidate trigger for LHS only if a CPD anomaly precedes the PL "
                    "hotspot in time AND co-locates with it in space, in the same ROI. Otherwise observed CPD "
                    "shifts may be ordinary SPV/ion relaxation unrelated to segregation."
                ),
                directions=(
                    "State-writing protocol S0->S5: capture CPD time series + PL hotspot time series in same fiducial ROI.\n"
                    "Statistical test: P[t_CPD < t_PL | same ROI] > 0.5 + delta supports H_beta (delta from drift budget).\n"
                    "Probes: fast-KPFM 3D scan + PL hyperspectral + low-fluence baseline subtraction."
                ),
                metric="P[t_CPD < t_PL] > 0.5+d",
                metric_caption="CAUSAL PRECEDENCE PROBABILITY",
            ),
            dict(
                num="H_g.", sub="Reversibility Boundary",
                fig_label="FIG. 3", fig_ratio="ref",
                fig_caption="Reversible vs irreversible: PL recovery, I2 gas trap, sealed-vs-open delta.",
                problem=(
                    "Reversible halide redistribution and irreversible chemical degradation share many PL "
                    "signatures. They are separable only by jointly considering PL, XPS chemical state, gas "
                    "trapping, and sealed-vs-open aging differences."
                ),
                directions=(
                    "Bayes-factor inference: K = P(y|M_reversible)/P(y|M_irreversible).\n"
                    "Decision: log10 K > 1 strong evidence reversible; < -1 strong evidence irreversible; |log10 K| < 1 demands sealed-vs-open follow-up.\n"
                    "Probes: PL recovery + AR-XPS Pb0/I0/PbI2 + RGA P_I2(t) + sealed vs open chamber pair."
                ),
                metric="|log10 K| > 1",
                metric_caption="BAYES FACTOR DECISION THRESHOLD",
            ),
        ],
    ),
    # S7
    dict(
        title="6. Method: State-Locked Physics-Informed Correlative Inference (SL-PCI)",
        intro=(
            "SL-PCI couples four steps. (i) Write a defined LHS state with operando optics on a fiducial-marked "
            "ROI. (ii) Lock the state by rapid quench or twin-sample dose array so destructive probes read the "
            "intended history. (iii) Read endpoint chemistry and structure. (iv) Infer a single 5-D latent state "
            "Z through a physics-informed Bayesian model in which each probe's artifact channel A_m is explicit, "
            "not hidden."
        ),
        page="0 7",
        cols=[
            dict(
                num="I.", sub="State Writing + Locking",
                fig_label="FIG. 1", fig_ratio="ref",
                fig_caption="State timeline S0-S5 with PL peak, CPD, lock points and transfer drift.",
                problem=(
                    "Operando PL and KPFM resolve dynamics (ms-min), but destructive probes read static endpoints. "
                    "Without locking, transfer drift mixes two different states under one label."
                ),
                directions=(
                    "Define S0 pristine, S1 pre-onset, S2 onset, S3 saturated, S4 recovery, S5 irreversible from PL peak shift and FWHM.\n"
                    "Lock by low-T quench in inert atmosphere; record history H = {I, t, T, atm, delay, probe dose}.\n"
                    "Validate lock fidelity: drift Delta_drift < mechanism signal Delta_mech."
                ),
                metric="Drift < signal",
                metric_caption="STATE-LOCK FIDELITY METRIC",
            ),
            dict(
                num="II.", sub="Correlative Coordinates",
                fig_label="FIG. 2", fig_ratio="ref",
                fig_caption="Fiducial Au/Cr grid; ROI mask categories; cross-platform registration.",
                problem=(
                    "PL hyperspectral, KPFM, TOF-SIMS, XPS and GIWAXS sit on different instruments with mismatched "
                    "spatial resolutions (300 nm / 50 nm / 100 nm / mm / mm). Pixel-to-pixel alignment fails; the "
                    "industry rarely admits this."
                ),
                directions=(
                    "Micron-scale Au/Cr fiducial grid printed before deposition.\n"
                    "ROI mask categories: grain center, GB, PL hotspot, high-CPD region, Br-rich region, passivated region.\n"
                    "Aggregate to ROI level; do not force pixel-to-pixel correspondence."
                ),
                metric=">=4 modalities co-registered",
                metric_caption="ROI-LEVEL CROSS-PLATFORM REGISTRATION",
            ),
            dict(
                num="III.", sub="Latent Inference",
                fig_label="FIG. 3", fig_ratio="ref",
                fig_caption="Bayesian graph: shared Z with per-probe likelihood H_m and artifact A_m.",
                problem=(
                    "Aggregating evidence with truth tables collapses uncertainty and conceals contradictions. "
                    "What is needed is a posterior over Z with proper covariance and identifiability analysis."
                ),
                directions=(
                    "Likelihood: y_m = H_m(Z; theta_m) + A_m + eta_m for each probe m.\n"
                    "Bayesian inference (HMC / VI) returns p(Z, theta | y) with covariance; Jacobian SVD reveals ill-conditioned directions before any sample is made.\n"
                    "Output: mechanism posteriors P(M | y) and Bayes factors between competing hypotheses."
                ),
                metric="cov(Z) + Bayes K",
                metric_caption="POSTERIOR DELIVERABLE PER SAMPLE",
            ),
        ],
    ),
    # S8
    dict(
        title="7. Forward Models Link Each Probe to One Shared Latent State Z",
        intro=(
            "Each characterization is a noisy projection of the same 5-D latent state Z = {x_Br/I, N_defect, "
            "phi_local, epsilon, C_deg}. Listing the forward model H_m(Z; theta_m) explicitly makes the inference "
            "reproducible and lets a Jacobian sensitivity analysis predict, before any experiment runs, which "
            "mechanism directions the chosen probe set can resolve and which require additional measurement."
        ),
        page="0 8",
        cols=[
            dict(
                num="I.", sub="Optical Projections",
                fig_label="FIG. 1", fig_ratio="ref",
                fig_caption="PL spectrum decomposed into Ruth-funneling-weighted phase components.",
                problem=(
                    "PL alone is ill-posed: peak position depends on x_Br/I, phase fraction phi, defect emission "
                    "and funneling l_diff simultaneously. PLQY and TRPL constrain N_defect via SRH lifetime."
                ),
                directions=(
                    "I_PL(E) = sum_p w_p * eta_p(N_defect) * L_p(E_g^p, sigma) * R(alpha, d), w_p propto V_D^-1 exp(-dEg/kT).\n"
                    "tau_eff^-1 = tau_SRH^-1(N_defect) + tau_rad^-1 + tau_surf^-1(S, d).\n"
                    "PLQY^-1 - 1 = 1/(k_rad*tau_bulk) + S_surf/(k_rad*d) separates bulk vs surface defects."
                ),
                metric="x_Br/I, phi, N_defect, e",
                metric_caption="RESOLVED LATENT COMPONENTS",
            ),
            dict(
                num="II.", sub="Electrostatic Projections",
                fig_label="FIG. 2", fig_ratio="ref",
                fig_caption="KPFM CPD = W0 - V_bb - V_SPV - V_ion - V_trap; time-frequency decomposition.",
                problem=(
                    "KPFM CPD is a 5-component convolution. Without time-frequency-intensity scanning plus a UPS "
                    "anchor for W0, no single mode can be attributed to local electric field."
                ),
                directions=(
                    "V_CPD = W0(chi_Br_surf) - V_bb(N_defect, N_d) - V_SPV(I, N_defect) - V_ion(t, mu_X) - V_trap(r, t).\n"
                    "Separate by time constants: V_SPV ~ ms (light step), V_trap ~ s (scan rate), V_ion ~ min (dark relax).\n"
                    "Anchor W0 with static UPS; remaining 4 components fit jointly with the latent inference."
                ),
                metric="chi_Br, N_defect, phi",
                metric_caption="RESOLVED LATENT COMPONENTS",
            ),
            dict(
                num="III.", sub="Chemical + Structural",
                fig_label="FIG. 3", fig_ratio="ref",
                fig_caption="AR-XPS depth integral; Vegard q(x_Br); Williamson-Hall strain vs size.",
                problem=(
                    "Structural FWHM is jointly broadened by strain epsilon and crystallite size L; chemical readouts "
                    "(XPS, SIMS) are perturbed by their own probes. Disentangling requires joint likelihood plus "
                    "dose-series controls."
                ),
                directions=(
                    "AR-XPS: I_Br(theta)/I_I(theta) = integral c_Br(z) exp(-z/lambda*cos theta) dz / integral c_I.\n"
                    "GIWAXS Vegard: q_hkl(x_Br) = q0 (1 - x_Br * Delta a/a); FWHM ~ epsilon*tan theta + K*lambda/L.\n"
                    "RGA: dN_I2/dt = R_loss(c_I, T, hv) * V_film, separates irreversible loss from reversible drift."
                ),
                metric="x_Br(z), e, L, C_deg",
                metric_caption="RESOLVED LATENT COMPONENTS",
            ),
        ],
    ),
    # S9
    dict(
        title="8. Identifiability Budget: Each Aim Closes Specific Mechanism Gaps",
        intro=(
            "Before any sample is made, the Jacobian J_mj = d H_m / d Z_j is built from all forward models. Its "
            "SVD reveals which mechanism directions the chosen probe set constrains well (large sigma_k) and which "
            "remain degenerate (small sigma_k). This identifiability budget makes Aim-by-Aim progress quantitative: "
            "each Aim adds probes that close specific ill-conditioned directions, monitored by condition number kappa."
        ),
        page="0 9",
        cols=[
            dict(
                num="I.", sub="Aim 1  (Year 1)",
                fig_label="FIG. 1", fig_ratio="ref",
                fig_caption="Synthetic-data pipeline (M0) + initial Jacobian condition number baseline.",
                problem=(
                    "Initial probe set (PL hyperspectral + UV-vis + static KPFM + UPS) cannot separate "
                    "N_defect_bulk from N_defect_surf; PLQY only constrains the sum. C_deg is unobservable."
                ),
                directions=(
                    "M0 (pre-sample): synthetic-data pipeline runs full Bayesian inference on known Z*, verifies recovery and computes kappa baseline.\n"
                    "Aim 1 probes: PL hyperspectral, UV-vis, static KPFM, UPS, baseline LHS state writing.\n"
                    "Resolves: x_Br_bulk, coarse phi. Unresolved: N_defect_bulk vs N_defect_surf, C_deg."
                ),
                metric="k baseline",
                metric_caption="INITIAL CONDITION NUMBER",
            ),
            dict(
                num="II.", sub="Aim 2  (Year 2)",
                fig_label="FIG. 2", fig_ratio="ref",
                fig_caption="kappa reduction after adding fast-KPFM 3D + TOF-SIMS + AR-XPS.",
                problem=(
                    "phi_local and N_defect_surf are not separable from optical data alone. Halide depth profile "
                    "x_Br/I(z) is unconstrained without surface-sensitive probes."
                ),
                directions=(
                    "Add fast-KPFM 3D scan (decomposes CPD 5 components), TOF-SIMS (z-profile), AR-XPS (surface chemistry), TRPL thickness gradient.\n"
                    "Resolves: + phi_local, N_defect_surf, x_Br/I(z).\n"
                    "Unresolved: epsilon vs L_crystallite still confounded (GIWAXS FWHM joint broadening)."
                ),
                metric="k <= 10^2",
                metric_caption="CONDITION NUMBER TARGET AFTER AIM 2",
            ),
            dict(
                num="III.", sub="Aim 3  (Year 3)",
                fig_label="FIG. 3", fig_ratio="ref",
                fig_caption="Full kappa ~ 10 with all 5 Z dimensions constrained; H_a/H_b/H_g testable.",
                problem=(
                    "Without separating epsilon from L, mechanism attributions for strain-driven LHS remain "
                    "ambiguous; C_deg requires explicit irreversible-chemistry probes."
                ),
                directions=(
                    "Add GIWAXS + Raman (epsilon vs L separation via Williamson-Hall + linewidth disorder), gas trapping/RGA (C_deg direct), sealed-vs-open aging.\n"
                    "Resolves: all 5 Z components.\n"
                    "Statistical hypotheses H_a / H_b / H_g become numerically testable with reported confidence intervals and Bayes factors."
                ),
                metric="k ~ 10",
                metric_caption="FULL IDENTIFIABILITY ACHIEVED",
            ),
        ],
    ),
    # S10
    dict(
        title="9. Three-Year Roadmap and Defensible Outputs",
        intro=(
            "The plan delivers two manuscripts plus an optional methods short paper, each anchored on falsifiable "
            "mechanism claims with reported uncertainty. Year 1 builds the pipeline on synthetic data plus a baseline "
            "sample platform. Year 2 publishes the SL-PCI proof-of-concept on the Br-rich vs passivation matrix. "
            "Year 3 closes the mechanism map across the high-Br stress regime and writes the thesis."
        ),
        page="1 0",
        cols=[
            dict(
                num="I.", sub="Year 1  (Foundation)",
                fig_label="FIG. 1", fig_ratio="ref",
                fig_caption="Reproducible sample + state platform; v0.1 Bayesian latent model.",
                problem=(
                    "No reproducible fiducial-grid platform exists for cross-platform same-area readout; locking-drift "
                    "budget is not quantified anywhere in the field."
                ),
                directions=(
                    "M0 pipeline: synthetic data + Bayesian inference + Jacobian SVD on PC, before any sample.\n"
                    "M1-M3: fiducial-grid WBG films (1.66-1.78 eV); baseline PL / UV-vis / XRD SOP.\n"
                    "M4-M9: define S0-S5 states; quantify low-T locking drift.\n"
                    "M10-M12: first latent model fit on real Aim-1 data."
                ),
                metric="drift < signal",
                metric_caption="LOCK FIDELITY VALIDATED",
            ),
            dict(
                num="II.", sub="Year 2  (Core Proof)",
                fig_label="FIG. 2", fig_ratio="ref",
                fig_caption="Manuscript 1: same-ROI PL-KPFM-TOF-SIMS with leave-one-modality-out CV.",
                problem=(
                    "Br-rich surface debate (Fang 2024 vs Navid 2026) has no decisive experiment; the field "
                    "lacks a same-ROI multimodal correlative test."
                ),
                directions=(
                    "M13-M18: same-ROI PL-KPFM-TOF-SIMS protocol on 2-composition x 3-treatment matrix.\n"
                    "M19-M21: artifact channel A_m calibration (dose series, sputter/transfer/delay).\n"
                    "M22-M24: cross-modal inference + leave-one-modality-out validation; write Manuscript 1.\n"
                    "Outcome: regression slope b1 = d_phi / d_chi_Br_surf with 95% CI tests H_alpha."
                ),
                metric="k <= 10^2; H_a tested",
                metric_caption="MANUSCRIPT 1 EVIDENCE CRITERION",
            ),
            dict(
                num="III.", sub="Year 3  (Mechanism Map)",
                fig_label="FIG. 3", fig_ratio="ref",
                fig_caption="Manuscript 2: mechanism regime diagram with Bayes-factor support.",
                problem=(
                    "No published mechanism regime diagram for WBG perovskite LHS specifies under which (light, "
                    "temperature, atmosphere) conditions each mechanism dominates with quantified uncertainty."
                ),
                directions=(
                    "M25-M30: add GIWAXS+Raman + XPS/RGA + sealed-vs-open; build regime diagram across Eg = 1.66-1.78 eV.\n"
                    "M31-M33: optional cryo-CL / cryo-EDS / Brillouin / NV-AFM cross-checks (only if SVD shows margin gain).\n"
                    "M34-M36: thesis integration + Manuscript 2: full mechanism contribution map with Bayes factors.\n"
                    "Defensible thesis claim: Br-rich surface effect is decided by the coupled latent state it creates, not by Br richness alone."
                ),
                metric="all 5 Z; log10 K",
                metric_caption="MANUSCRIPT 2 + THESIS COMPLETE",
            ),
        ],
    ),
]


def find_shape_by_text(slide, target_text, x_min=None, x_max=None, y_min=None, y_max=None):
    """Find a shape whose text frame's first non-empty paragraph matches target_text (or starts with it)."""
    target_clean = target_text.strip()
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        try:
            lx = Emu(sh.left).inches if sh.left is not None else None
            ly = Emu(sh.top).inches if sh.top is not None else None
        except Exception:
            lx, ly = None, None
        if x_min is not None and lx is not None and lx < x_min - 0.05:
            continue
        if x_max is not None and lx is not None and lx > x_max + 0.05:
            continue
        if y_min is not None and ly is not None and ly < y_min - 0.05:
            continue
        if y_max is not None and ly is not None and ly > y_max + 0.05:
            continue
        # Concatenated text from all paragraphs
        all_text = " | ".join((p.text or "").strip() for p in sh.text_frame.paragraphs if (p.text or "").strip())
        if all_text.startswith(target_clean) or target_clean in all_text:
            return sh
    return None


def set_shape_text(shape, new_text, lines=None):
    """Replace shape text while preserving the first run's font formatting.

    If `lines` is a list, populate paragraphs accordingly (one paragraph per line).
    Otherwise treat new_text as a single string; '\n' splits into paragraphs.
    """
    tf = shape.text_frame
    # Capture the first run's font reference (a deep copy) before clearing.
    template_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            template_run = r
            break
        if template_run is not None:
            break

    # Determine list of paragraph strings
    if lines is None:
        lines = new_text.split("\n")
    lines = [l for l in lines]

    # Wipe existing paragraphs except the first
    txBody = tf._txBody
    pPr_first = None
    # Remove all existing <a:p> children
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for p in list(txBody.findall("a:p", nsmap)):
        txBody.remove(p)

    # Add a fresh <a:p> per line, copying template run formatting
    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for i, line in enumerate(lines):
        p_el = etree.SubElement(txBody, f"{a_ns}p")
        if template_run is not None:
            new_run = copy.deepcopy(template_run._r)
            # Replace text content in the new run
            for t in new_run.findall(f"{a_ns}t"):
                new_run.remove(t)
            t_el = etree.SubElement(new_run, f"{a_ns}t")
            t_el.text = line
            p_el.append(new_run)
        else:
            r_el = etree.SubElement(p_el, f"{a_ns}r")
            t_el = etree.SubElement(r_el, f"{a_ns}t")
            t_el.text = line


def duplicate_slide(prs, source_idx):
    """Duplicate slide at source_idx and append to end of presentation. Returns new slide."""
    src = prs.slides[source_idx]
    blank_layout = src.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    # Remove all shapes from new_slide first (the layout placeholders)
    spTree = new_slide.shapes._spTree
    for sh in list(new_slide.shapes):
        sp = sh._element
        sp.getparent().remove(sp)

    # Copy all shapes from src
    src_spTree = src.shapes._spTree
    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    p_ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    for child in src_spTree:
        tag = etree.QName(child.tag).localname
        if tag in ("nvGrpSpPr", "grpSpPr"):
            continue
        new_child = copy.deepcopy(child)
        spTree.append(new_child)

    # Copy rels for images
    from pptx.parts.image import ImagePart
    for rel in src.part.rels.values():
        if "image" in rel.reltype or "chart" in rel.reltype or "hyperlink" in rel.reltype:
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    return new_slide


def fill_slide(slide, content):
    """Fill new slide (cloned from S4) with the SL-PCI content dict."""
    # Title
    sh = find_shape_by_text(slide, "3. Three Critical Challenges", y_max=0.3)
    if sh is None:
        # fallback: find shape at L=0.07 T=0.1
        for s in slide.shapes:
            if s.has_text_frame and s.left is not None and Emu(s.top).inches < 0.25 and Emu(s.left).inches < 0.5:
                if s.width is not None and Emu(s.width).inches > 5:
                    sh = s
                    break
    if sh is not None:
        set_shape_text(sh, content["title"])

    # Intro paragraph (T ~ 0.62, wide box L ~ 0.18)
    for s in slide.shapes:
        if s.has_text_frame and s.top is not None:
            try:
                if 0.55 < Emu(s.top).inches < 0.75 and Emu(s.width).inches > 10:
                    set_shape_text(s, content["intro"])
                    break
            except Exception:
                continue

    # Page number (T ~ 7.32, L ~ 12.78)
    for s in slide.shapes:
        if s.has_text_frame and s.top is not None and s.left is not None:
            try:
                if Emu(s.top).inches > 7.2 and Emu(s.left).inches > 12.5:
                    set_shape_text(s, content["page"])
                    break
            except Exception:
                continue

    # Three columns: identify by approximate X range
    col_xranges = [(0.0, 4.5), (4.5, 8.7), (8.7, 13.4)]
    col_data = content["cols"]

    for col_idx, (xmin, xmax) in enumerate(col_xranges):
        col = col_data[col_idx]

        # Column number (I./II./III.) at T ~ 1.28, small box
        for s in slide.shapes:
            if not s.has_text_frame: continue
            try:
                lx, ly = Emu(s.left).inches, Emu(s.top).inches
                w = Emu(s.width).inches
            except Exception: continue
            if xmin <= lx < xmax and 1.2 < ly < 1.32 and w < 0.6:
                txt = " ".join((p.text or "").strip() for p in s.text_frame.paragraphs).strip()
                if txt in ("I.", "II.", "III."):
                    set_shape_text(s, col["num"])
                    break

        # Column sub-title at T ~ 1.34 (Stability/Upscaling/Efficiency)
        for s in slide.shapes:
            if not s.has_text_frame: continue
            try:
                lx, ly = Emu(s.left).inches, Emu(s.top).inches
                w = Emu(s.width).inches
            except Exception: continue
            if xmin <= lx < xmax and 1.3 < ly < 1.4 and 0.6 < w < 1.5:
                txt = " ".join((p.text or "").strip() for p in s.text_frame.paragraphs).strip()
                if txt in ("Stability", "Upscaling", "Efficiency"):
                    set_shape_text(s, col["sub"])
                    break

        # Fig label (FIG. n) at T ~ 1.87
        for s in slide.shapes:
            if not s.has_text_frame: continue
            try:
                lx, ly = Emu(s.left).inches, Emu(s.top).inches
            except Exception: continue
            if xmin <= lx < xmax and 1.82 < ly < 1.92:
                txt = " ".join((p.text or "").strip() for p in s.text_frame.paragraphs).strip()
                if txt.startswith("FIG."):
                    set_shape_text(s, col["fig_label"])

        # Suggested fig caption at T ~ 2.87
        for s in slide.shapes:
            if not s.has_text_frame: continue
            try:
                lx, ly = Emu(s.left).inches, Emu(s.top).inches
            except Exception: continue
            if xmin <= lx < xmax and 2.8 < ly < 2.95:
                txt = " ".join((p.text or "").strip() for p in s.text_frame.paragraphs).strip()
                if txt.startswith("Suggested:"):
                    set_shape_text(s, col["fig_caption"])
                    break

        # PROBLEM body at T ~ 4.44
        for s in slide.shapes:
            if not s.has_text_frame: continue
            try:
                lx, ly = Emu(s.left).inches, Emu(s.top).inches
                w = Emu(s.width).inches
            except Exception: continue
            if xmin <= lx < xmax and 4.35 < ly < 4.55 and w > 3.0:
                set_shape_text(s, col["problem"])
                break

        # DIRECTIONS body at T ~ 5.32 (multi-line)
        for s in slide.shapes:
            if not s.has_text_frame: continue
            try:
                lx, ly = Emu(s.left).inches, Emu(s.top).inches
                w = Emu(s.width).inches
            except Exception: continue
            if xmin <= lx < xmax and 5.25 < ly < 5.45 and w > 3.0:
                lines = col["directions"].split("\n")
                set_shape_text(s, "", lines=lines)
                break

        # Bottom metric value at T ~ 6.17-6.32 (small width)
        for s in slide.shapes:
            if not s.has_text_frame: continue
            try:
                lx, ly = Emu(s.left).inches, Emu(s.top).inches
                w = Emu(s.width).inches
            except Exception: continue
            if xmin <= lx < xmax and 6.15 < ly < 6.35 and w < 1.2:
                txt = " ".join((p.text or "").strip() for p in s.text_frame.paragraphs).strip()
                if txt in ("> 1000 h", ">= 800 cm²", "≥ 800 cm²", "26.7 %"):
                    set_shape_text(s, col["metric"])
                    break

        # Bottom metric caption at T ~ 6.28-6.44 (wider)
        for s in slide.shapes:
            if not s.has_text_frame: continue
            try:
                lx, ly = Emu(s.left).inches, Emu(s.top).inches
                w = Emu(s.width).inches
            except Exception: continue
            if xmin <= lx < xmax and 6.25 < ly < 6.5 and w > 0.8:
                txt = " ".join((p.text or "").strip() for p in s.text_frame.paragraphs).strip()
                if "OPERATIONAL TARGET" in txt or "MODULE-SCALE RESEARCH TARGET" in txt or "SINGLE-JUNCTION CERTIFIED RECORD" in txt:
                    set_shape_text(s, col["metric_caption"])
                    break


def remove_slide(prs, idx):
    """Remove slide at idx from sldIdLst (correctly drops the relationship)."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    target = slides_list[idx]
    # also drop the rel from package
    rId = target.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
    prs.part.drop_rel(rId)
    xml_slides.remove(target)


def main():
    prs = Presentation(str(SRC))
    n_slides_initial = len(prs.slides)
    print(f"Initial slides: {n_slides_initial}")

    # Step 1: remove slides 5-10 (indices 4-9, six slides)
    # We must remove from the highest index downward.
    for idx in [9, 8, 7, 6, 5, 4]:
        if idx < len(prs.slides):
            print(f"  removing slide {idx+1}")
            remove_slide(prs, idx)

    print(f"After removal: {len(prs.slides)}")

    # Step 2: duplicate slide 4 (index 3) six times for new S5-S10
    new_slides = []
    for i in range(6):
        ns = duplicate_slide(prs, 3)
        new_slides.append(ns)
        print(f"  duplicated -> new slide {len(prs.slides)}")

    # Step 3: but newly added slides are at end of deck; we need them at positions 4-9
    # (between slide 4 and what is currently slide 5, originally slide 11 in the source).
    # Reorder sldIdLst.
    xml_slides = prs.slides._sldIdLst
    sld_list = list(xml_slides)
    # Current ordering: [s1, s2, s3, s4, s11', s12', s13', s14', s15', ..., s19_post, NEW1..NEW6]
    # We want:           [s1, s2, s3, s4, NEW1..NEW6, s11', s12', ...]
    # The newly added 6 are the last 6 in current ordering.
    head = sld_list[:4]          # s1-s4
    new6 = sld_list[-6:]
    middle = sld_list[4:-6]
    new_order = head + new6 + middle
    # Re-append in new order
    for el in sld_list:
        xml_slides.remove(el)
    for el in new_order:
        xml_slides.append(el)

    print(f"Final slide count: {len(prs.slides)}")
    print(f"Reordered: head=4, new=6, rest={len(middle)}")

    # Step 4: fill in new slides 5-10 with SL-PCI content
    for i, content in enumerate(SLIDES):
        slide = prs.slides[4 + i]
        print(f"Filling new slide {5 + i}: {content['title'][:60]}")
        fill_slide(slide, content)

    # Step 5: save
    prs.save(str(OUT))
    print(f"\nSaved: {OUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
